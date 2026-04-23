"""
Document Generation Tools — Create docx, xlsx, pptx, and PDF files.
Mirrors Cowork's document creation skills.
"""

import os
from langchain_core.tools import tool
from src.tools.workspace_context import tool_workspace_root


def _out_path(filename: str) -> str:
    """Resolve output path inside the active workspace."""
    root = tool_workspace_root()
    path = os.path.join(root, filename)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    return path


@tool
def create_docx(filename: str, content: str, title: str = "") -> str:
    """
    Creates a Word document (.docx) in the workspace.

    Content supports simple markdown-like formatting:
    - Lines starting with # become headings
    - Lines starting with - become bullet points
    - Other lines become normal paragraphs
    - Empty lines create paragraph breaks

    Args:
        filename: Output filename (e.g. 'report.docx').
        content: The document content with simple formatting.
        title: Optional document title (added as heading at top).
    """
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph("")
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped.startswith("1. ") or stripped.startswith("2. ") or stripped.startswith("3. "):
            doc.add_paragraph(stripped[3:], style="List Number")
        else:
            doc.add_paragraph(stripped)

    if not filename.endswith(".docx"):
        filename += ".docx"
    path = _out_path(filename)
    doc.save(path)
    return f"✅ Created Word document: {filename}"


@tool
def create_xlsx(filename: str, data: str, sheet_name: str = "Sheet1") -> str:
    """
    Creates an Excel spreadsheet (.xlsx) in the workspace.

    Data format: CSV-like text where each line is a row and values are comma-separated.
    The first row is treated as headers.

    Example data:
        Name, Age, City
        Alice, 30, Bangkok
        Bob, 25, Tokyo

    Args:
        filename: Output filename (e.g. 'data.xlsx').
        data: CSV-like text content for the spreadsheet.
        sheet_name: Name of the worksheet.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    rows = [line.strip() for line in data.strip().split("\n") if line.strip()]
    for i, row_text in enumerate(rows):
        cells = [c.strip() for c in row_text.split(",")]
        for j, val in enumerate(cells, 1):
            cell = ws.cell(row=i + 1, column=j, value=val)
            if i == 0:
                cell.font = Font(bold=True)

    if not filename.endswith(".xlsx"):
        filename += ".xlsx"
    path = _out_path(filename)
    wb.save(path)
    return f"✅ Created spreadsheet: {filename} ({len(rows)} rows)"


@tool
def create_pptx(filename: str, slides_content: str, title: str = "") -> str:
    """
    Creates a PowerPoint presentation (.pptx) in the workspace.

    Slides are separated by '---' on its own line.
    First line of each slide becomes the title, rest becomes bullet points.

    Example:
        Introduction
        - Welcome to the presentation
        - Overview of topics
        ---
        Key Findings
        - Finding 1: Revenue grew 20%
        - Finding 2: User base doubled

    Args:
        filename: Output filename (e.g. 'presentation.pptx').
        slides_content: Slide content separated by '---'.
        title: Optional presentation title for the first slide.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation()

    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title

    raw_slides = slides_content.split("---")
    for raw in raw_slides:
        lines = [l.strip() for l in raw.strip().split("\n") if l.strip()]
        if not lines:
            continue

        slide_title = lines[0].lstrip("# ")
        bullets = []
        for line in lines[1:]:
            if line.startswith("- ") or line.startswith("* "):
                bullets.append(line[2:])
            else:
                bullets.append(line)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        if bullets:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b

    if not filename.endswith(".pptx"):
        filename += ".pptx"
    path = _out_path(filename)
    prs.save(path)
    slide_count = len(prs.slides)
    return f"✅ Created presentation: {filename} ({slide_count} slides)"


@tool
def create_pdf(filename: str, content: str, title: str = "") -> str:
    """
    Creates a PDF document in the workspace from text content.

    Args:
        filename: Output filename (e.g. 'report.pdf').
        content: Text content for the PDF.
        title: Optional title at the top of the document.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "Error: PyMuPDF not installed. Run: pip install pymupdf"

    doc = fitz.open()
    page = doc.new_page()
    y = 72  # top margin

    if title:
        page.insert_text((72, y), title, fontsize=18, fontname="helv")
        y += 30

    for line in content.split("\n"):
        if y > page.rect.height - 72:
            page = doc.new_page()
            y = 72
        stripped = line.strip()
        if stripped.startswith("# "):
            page.insert_text((72, y), stripped[2:], fontsize=14, fontname="helv")
            y += 22
        else:
            page.insert_text((72, y), stripped if stripped else " ", fontsize=10, fontname="helv")
            y += 14

    if not filename.endswith(".pdf"):
        filename += ".pdf"
    path = _out_path(filename)
    doc.save(path)
    doc.close()
    return f"✅ Created PDF: {filename}"
